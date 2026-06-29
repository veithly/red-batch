"""Add the live UiPath Orchestrator proof screenshots to the pitch deck.

Slide 4 ("Solution architecture") has an empty bottom half. Place two live
screenshots there side by side -- the coded agent published (Tenant -> Packages)
and deployed (Shared -> Processes) -- with a one-line caption. Preserves the
rest of the template design. Run: python3 scripts/add_deck_proof.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DECK = "docs/red-batch-deck.pptx"
IMG_PKG = "docs/uipath-orchestrator-package.png"
IMG_PROC = "docs/uipath-orchestrator-process.png"
BODY = RGBColor(0x1D, 0x1D, 0x1F)
BORDER = RGBColor(0xC8, 0xCC, 0xD0)

# Slide-4 architecture text ends at ~4.19in; place the proof below it.
CAP_TOP = 4.34
IMG_TOP = 4.66
IMG_W = 4.8
LEFT_L = 1.5
RIGHT_L = 7.02

prs = Presentation(DECK)
slide = list(prs.slides)[3]  # slide 4

cap = slide.shapes.add_textbox(Inches(0.36), Inches(CAP_TOP), Inches(12.6), Inches(0.30))
tf = cap.text_frame
tf.word_wrap = True
para = tf.paragraphs[0]
para.alignment = PP_ALIGN.CENTER
r = para.add_run()
r.text = ("Live in the UiPath tenant \u2014 coded agent published (Tenant \u2192 Packages, left) "
          "and deployed as a Process (Shared \u2192 Processes, right)")
r.font.size = Pt(12)
r.font.bold = True
r.font.color.rgb = BODY

for path, left in ((IMG_PKG, LEFT_L), (IMG_PROC, RIGHT_L)):
    pic = slide.shapes.add_picture(path, Inches(left), Inches(IMG_TOP), width=Inches(IMG_W))
    pic.line.color.rgb = BORDER
    pic.line.width = Pt(0.75)

prs.save(DECK)
print("SAVED", DECK)
